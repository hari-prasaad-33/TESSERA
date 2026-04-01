#!/usr/bin/env python3
"""
Convert Tessera One pitch deck Markdown to PowerPoint (.pptx).
Requires: pip install python-pptx

Usage: python pitch_to_pptx.py [input.md] [output.pptx]
"""

import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    sys.exit(1)


def parse_slides(md_path: Path) -> list[dict]:
    """Parse markdown into list of {title, content, is_table} dicts."""
    text = md_path.read_text(encoding="utf-8")

    # Split by ## Slide N: pattern
    slide_pattern = re.compile(r"^## Slide \d+:\s*(.+)$", re.MULTILINE)
    parts = re.split(slide_pattern, text)

    # First part is preamble (title, intro); rest are alternating (title, content)
    slides = []
    for i in range(1, len(parts), 2):
        title = parts[i].strip()
        content = parts[i + 1].strip() if i + 1 < len(parts) else ""
        # Remove leading --- and trailing ---
        content = re.sub(r"^---+|---+$", "", content).strip()

        # Check for table
        is_table = "|" in content and content.strip().startswith("|")

        slides.append({"title": title, "content": content, "is_table": is_table})

    return slides


def add_text_slide(prs, title: str, content: str):
    """Add a slide with title and bullet content."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(0.8)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x1a, 0x1a, 0x1a)

    # Content
    top = Inches(1.3)
    height = Inches(5.5)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    lines = content.strip().split("\n")

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Strip markdown bold/italic
        line = re.sub(r"\*\*(.+?)\*\*", r"\1", line)
        line = re.sub(r"\*(.+?)\*", r"\1", line)

        # Tables: add as plain text with spacing
        if line.startswith("|") and "---" not in line:
            line = "  ".join(c.strip() for c in line.split("|")[1:-1])
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
            p.space_before = Pt(2)
            continue

        # Skip table separator row
        if re.match(r"^\|[\s\-:]+\|", line):
            continue

        p = tf.add_paragraph()
        # Bullet or numbered
        if line.startswith("- "):
            p.text = "• " + line[2:]
            p.level = 0
        elif re.match(r"^\d+\.\s", line):
            p.text = line
            p.level = 0
        else:
            p.text = line
            p.level = 0

        p.font.size = Pt(14)
        p.space_before = Pt(6)


def add_table_slide(prs, title: str, content: str):
    """Add a slide with title and a table."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(0.8)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x1a, 0x1a, 0x1a)

    # Parse table
    rows = []
    for line in content.strip().split("\n"):
        if line.strip().startswith("|") and "---" not in line:
            cells = [c.strip() for c in line.split("|")[1:-1]]
            rows.append(cells)

    if not rows:
        add_text_slide(prs, title, content)
        return

    num_rows = len(rows)
    num_cols = max(len(r) for r in rows)

    # Create table
    col_width = Inches(2.0)
    table_width = Inches(2.0 * num_cols)
    table_height = Inches(0.35 * num_rows)
    table_left = Inches(0.5)
    table_top = Inches(1.5)

    table = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height).table

    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            if j < num_cols:
                cell = table.cell(i, j)
                cell.text = cell_text
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                if i == 0:
                    cell.text_frame.paragraphs[0].font.bold = True

    # Add any content after the table
    after_table = []
    in_table = True
    for line in content.strip().split("\n"):
        if line.strip().startswith("|"):
            if "---" in line:
                in_table = False
            continue
        if not in_table and line.strip():
            after_table.append(line)

    if after_table:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.95 + 0.35 * num_rows), Inches(9), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for line in after_table:
            line = re.sub(r"\*\*(.+?)\*\*", r"\1", line)
            line = re.sub(r"\*(.+?)\*", r"\1", line)
            p = tf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(12)
            p.space_before = Pt(4)


def main():
    script_dir = Path(__file__).parent
    md_path = Path(script_dir) / "TESSERA_ONE_PITCH_DECK.md"
    out_path = Path(script_dir) / "Tessera_One_Pitch_Deck.pptx"

    if len(sys.argv) >= 2:
        md_path = Path(sys.argv[1])
    if len(sys.argv) >= 3:
        out_path = Path(sys.argv[2])

    if not md_path.exists():
        print(f"Input file not found: {md_path}")
        sys.exit(1)

    slides_data = parse_slides(md_path)
    if not slides_data:
        print("No slides found in input.")
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for s in slides_data:
        if s["is_table"]:
            add_table_slide(prs, s["title"], s["content"])
        else:
            add_text_slide(prs, s["title"], s["content"])

    prs.save(out_path)
    print(f"Created: {out_path}")


if __name__ == "__main__":
    main()
